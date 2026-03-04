from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from copy import deepcopy
from io import BytesIO
from pptx import Presentation
from copy import deepcopy
from io import BytesIO


BASE_DIR = Path(__file__).resolve().parent
OBS = BASE_DIR / "Observations"
OBS.mkdir(exist_ok=True)

def paste_df_to_ppt(slide, df, left_inch, top_inch):

    # Convert position
    left = Inches(left_inch)
    top = Inches(top_inch)

    rows = df.shape[0] + 1
    cols = df.shape[1]

    # Create table (temporary width, we'll auto-size columns after)
    table = slide.shapes.add_table(
        rows, cols,
        left, top,
        Inches(1), Inches(4)
    ).table

    # ---------- ROW HEIGHT ----------
    for r in range(rows):
        table.rows[r].height = Pt(14)

    # ---------- HEADER ----------
    header_color = RGBColor(22, 62, 100)  # dark blue

    for col_idx, col_name in enumerate(df.columns):
        cell = table.cell(0, col_idx)
        cell.text = str(col_name)

        p = cell.text_frame.paragraphs[0]
        p.font.bold = True
        p.font.size = Pt(11)
        p.font.color.rgb = RGBColor(255, 255, 255)

        cell.fill.solid()
        cell.fill.fore_color.rgb = header_color

    # ---------- BODY ----------
    for row_idx, row in enumerate(df.itertuples(), start=1):
        row_values = row[1:]

        bold_flag = any(
            "total" in str(v).lower() and ":" in str(v)
            for v in row_values
        ) or any(
            x.lower() in str(v).lower()
            for v in row_values
            for x in ["cyclical exposure", "sensitive exposure", "defensive exposure"]
        )

        for col_idx, value in enumerate(row_values):
            cell = table.cell(row_idx, col_idx)
            cell.text = str(value)

            p = cell.text_frame.paragraphs[0]
            p.font.size = Pt(9)
            p.font.color.rgb = RGBColor(0, 0, 0)
            p.font.bold = bold_flag

            cell.fill.solid()
            cell.fill.fore_color.rgb = RGBColor(255, 255, 255)

    # ---------- AUTO COLUMN WIDTH ----------
    for col_idx in range(cols):

        max_width = Inches(0.6)

        # Header width
        header_text = str(df.columns[col_idx])
        header_width = Inches(len(header_text) * 0.085 + 0.3)
        max_width = max(max_width, header_width)

        # Body width
        for row_idx in range(len(df)):
            cell_text = str(df.iloc[row_idx, col_idx])
            cell_width = Inches(len(cell_text) * 0.065 + 0.3)
            max_width = max(max_width, cell_width)

        table.columns[col_idx].width = max_width



def build_master_report(results):

    prs = Presentation(OBS / "MasterTemplate.pptx")

    sections = {
        "Taxable Bonds in Non-Qualified Accounts": [
            {"df": results["taxable_bonds_in_nq"], "left": 0.2, "top": 3}
        ],
        "Tax Cost Ratio": [
            {"df": results["tax_cost_ratio"], "left": 0.2, "top": 2.6}
        ],
        "Expense Ratio": [
            {"df": results["high_expense_ratio"], "left": 0.2, "top": 2.5}
        ],
        "Muni Bonds in Qualified Accounts": [
            {"df": results["muni_bond_qualified"], "left": 0.2, "top": 3.1}
        ],
        "Active Passive Management": [
            {"df": results["active_passive"], "left": 0.2, "top": 2.5}
        ],
        "Tracking Error": [
            {"df": results["tracking_error"], "left": 0.2, "top": 2.5}
        ],
        "Specific Exposure": [
            {"df": results["countries_sector_rows"], "left": 0.2, "top": 2.5}
        ],
        "Inverse/ Leveraged Funds": [
            {"df": results["leverage_inverse_df"], "left": 0.2, "top": 2.5}
        ],
        "Digital Assets": [
            {"df": results["digital_df"], "left": 0.2, "top": 2.7}
        ],
        "Retail Share Class": [
            {"df": results["retail_class"], "left": 0.2, "top": 2.5}
        ],
        "Target Date Funds": [
            {"df": results["target_date_fund"], "left": 0.2, "top": 2.5}
        ],
        "Current Allocation By Sector": [
            {"df": results["sectors"], "left": 0.1, "top": 1.3}
        ],
        "Individual Stocks": [
            {"df": results["individual_stocks_df"], "left": 0.2, "top": 2.9}
        ],
        "Potential Conflict Of Interest": [
            {"df": results["fund_families"], "left": 0.2, "top": 2.5}
        ],

        "Underperforming Funds": [
            {"df": results["dilution"], "left": 0.2, "top": 2.5}
        ],


        "Global Asset Allocation": [
            {"df": results["summary_gaa"], "left": 0.1, "top": 0.7},
            {"df": results["general_gaa"], "left": 0.1, "top": 3.99}
        ],


        "Qualified Accounts:": [
            {"df": results["q_gaa"], "left": 0.1, "top": 0.7},
            {"df": results["nq_gaa"], "left": 0.1, "top": 3.99}
        ],


        "Portfolio Overview": [
            {"df": results["security_type_table"], "left": 0.2, "top": 0.8},
            {"df": results["performance_summary"], "left": 0.2, "top": 2.8},
            {"df": results["expense_summary"], "left": 4.5, "top": 0.8},
            {"df": results["aggregate_stylebox"], "left": 4.5, "top": 2.8},
            {"df": results["mini_fixed_income"], "left": 4.5, "top": 4.5}
        ],



        "Fixed Income Exposure": [
            {"df": results["credit_yield"], "left": 0.2, "top": 2.8},
            {"df": results["credit_duration"], "left": 0.2, "top": 4.5},
            {"df": results["credit_issuer_exposure"], "left": 4.5, "top": 2.8},
            {"df": results["credit_quality_exposure"], "left": 4.5, "top": 4.5}
        ],


    }

    slides_to_delete = []

    for i, slide in enumerate(prs.slides):

        # --- Detect title safely ---
        title = None
        candidates = []

        for shape in slide.shapes:
            if shape.has_text_frame:
                text = shape.text.strip()
                if text:
                    candidates.append(text)

        if candidates:
            # Shortest text is almost always the section title
            candidates_sorted = sorted(candidates, key=len)
            title = candidates_sorted[0]

        if not title:
            continue

        # --- Match title case-insensitive ---
        matched_key = next(
            (key for key in sections if key.lower().strip() == title.lower().strip()),
            None
        )

        if not matched_key:
            continue

        blocks = sections[matched_key]

        # --- Check if all DataFrames empty ---
        all_empty = all(
            block["df"] is None or block["df"].empty
            for block in blocks
        )

        if all_empty:
            slides_to_delete.append(i)
            continue

        # --- Paste tables ---
        for block in blocks:
            df = block["df"]
            if df is not None and not df.empty:
                paste_df_to_ppt(
                    slide,
                    df,
                    block["left"],
                    block["top"]
                )

    # --- Delete slides in reverse order ---
    for index in reversed(slides_to_delete):
        rId = prs.slides._sldIdLst[index].rId
        prs.part.drop_rel(rId)
        del prs.slides._sldIdLst[index]

    buffer = BytesIO()
    prs.save(buffer)
    buffer.seek(0)

    return buffer